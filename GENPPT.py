import json
import os
from pptx import Presentation
from functions import *


json_file = os.path.join("data.json")
with open(json_file, 'r') as file:
    insight_data = json.load(file)
print(insight_data)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

for insight_key, insight_value in insight_data.items():
    title = insight_value.get('title', '')
    description = insight_value.get('description', '')
    require_chart = insight_value.get('requireChart', 'false').lower() == 'true'

    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)

    prs = add_details_slide(prs, slide, title, description)

    if require_chart:
        chart_details = insight_value.get('details', {})
        chart_type = chart_details['chart']['type']
        title = chart_details['title']['text']
        categories = chart_details['xAxis']['categories']
        # x_label = chart_details['xAxis']
        # y_label = chart_details['yAxis']['title']['text']
        color_list = [RGBColor(0, 178, 255), RGBColor(0, 229, 124), RGBColor(95, 139, 188), RGBColor(81, 79, 196)]
        data_series = chart_details['series']
        if chart_type == 'line':
            # y_label = chart_details['yAxis']['title']['text']
            y_label =''
            prs = plot_line_chart(prs, slide, categories, data_series, title,y_label, color_list)
        elif chart_type == 'column':
            # y_label = chart_details['yAxis']['title']['text']
            y_label =''
            prs = plot_column_chart(prs, slide, categories, data_series, title,y_label, color_list)
        # elif chart_type == 'heatmap':
        #     prs = plot_heatmap_chart(prs,slide, categories, data_series, title)
        elif chart_type == 'pie':
            y_label = chart_details['yAxis']['title']['text']
            prs = plot_pie_chart(prs, slide, categories, data_series, title, color_list)

        # elif chart_type == 'area':
        #     prs = plot_area_chart(prs, slide, categories, data_series, title,color_list)
        # elif chart_type == 'scatter':
        #    prs = plot_scatter_chart(prs, slide, categories, data_series, title)

output_file = os.path.join('Output', 'Insights_Presentation_with_Charts.pptx')
prs.save(output_file)


#ppt_function.py
 
from pptx.chart.data import CategoryChartData,XyChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_LEGEND_POSITION, XL_LABEL_POSITION,XL_MARKER_STYLE,XL_TICK_MARK
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_details_slide(prs, slide, title, description):

    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    content_placeholder = slide.placeholders[1]
    content_placeholder.text = description
    content_shape = content_placeholder.text_frame
    content_paragraph = content_shape.paragraphs[0]
    content_paragraph.alignment = PP_ALIGN.LEFT
    content_font = content_paragraph.font
    content_font.size = Pt(20)
    return prs
def plot_line_chart(prs, slide,categories, data_series,title,y_label,color_list):

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series in data_series:
        chart_series = chart_data.add_series(series['name'])
        for value in series['data']:
            chart_series.add_data_point(value)

    x, y, cx, cy = Inches(9), Inches(2), Inches(7), Inches(6)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    for i, series in enumerate(chart.plots[0].series):
        series.format.line.color.rgb = color_list[i]
        series.has_markers = True
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = 7
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color_list[i]
        series.marker.format.line.color.rgb = color_list[i]

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.axis_title.text_frame.text = y_label
    # value_axis.tick_labels.number_format = '0,,"M"'
    chart.value_axis.tick_labels.font.size = Pt(10)

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.rotation = 45
    # chart.category_axis.axis_title.text_frame.text = y_label
    plot = chart.plots[0]
    plot.has_data_labels = False

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(8)

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    chart.value_axis.major_gridlines.format.line.width = 0
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(255, 255, 255)

    return prs

def plot_scatter_chart(prs, slide,categories, data_series,title):

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series in data_series:
        chart_series = chart_data.add_series(series['name'])
        for value in series['data']:
            chart_series.add_data_point(value)

    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
    ).chart

    for i, series in enumerate(chart.plots[0].series):
        # line_color = hex_to_rgb(colors[i % len(colors)])
        series.format.line.color.rgb = RGBColor(0, 137, 253)

        series.has_markers = True
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = 7

        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = RGBColor(0, 137, 253)

        series.marker.format.line.color.rgb = RGBColor(0, 137, 253)

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.rotation = 45

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(8)
    # category_axis.axis_title.text_frame.text = categories  # X-axis label
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = False
    value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.tick_labels.font.size = Pt(10)
    # value_axis.axis_title.text_frame.text = title
    # value_axis.tick_labels.number_format = '0,,"M"'
    chart.value_axis.tick_labels.font.size = Pt(10)

    plot = chart.plots[0]
    plot.has_data_labels = False

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(8)

    # Configure the chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    chart.value_axis.major_gridlines.format.line.width = 0
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(255, 255, 255)

    return prs

def plot_column_chart(prs, slide, categories, data_series, title,y_label,color_list):
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series in data_series:
        chart_series = chart_data.add_series(series['name'])
        for i, value in enumerate(series['data']):
            chart_series.add_data_point(value, categories[i])

    x, y, cx, cy = Inches(9), Inches(2), Inches(7), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    for series in chart.series:
        for point in series.points:
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 178, 255)



    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = False

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.has_major_gridlines = False
    # chart.value_axis.axis_title.text_frame.text = y_label
    chart.value_axis.axis_title.text_frame.text = y_label
    return prs
def plot_pie_chart(prs, slide, categories, data_series, title,color_list):
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for i,series in enumerate(data_series):
        chart_series = chart_data.add_series(series['name'])
        for value in series['data']:
            chart_series.add_data_point(value,categories[i])

    x, y, cx, cy = Inches(8), Inches(3), Inches(7), Inches(5)

    # Create the pie chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    for series in chart.series:
        for i, point in enumerate(series.points):
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = color_list[i]

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = True  # Add a legend
    chart.plots[0].has_data_labels = True  # Show data labels on the pie slices

    return prs

def plot_area_chart(prs, slide, categories, data_series, title,color_list):
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series in data_series:
        chart_series = chart_data.add_series(series['name'])
        for value in series['data']:
            chart_series.add_data_point(value)

    x, y, cx, cy = Inches(8), Inches(2), Inches(7), Inches(5)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA, x, y, cx, cy, chart_data
    ).chart

    for i,series in enumerate(chart.series):
        for point in series.points:
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = color_list[i]

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
    chart.value_axis.tick_labels.font.size = Pt(10)
    # chart.value_axis.axis_title.text_frame.text = y_label
    # value_axis.tick_labels.number_format = '0,,"M"'
    chart.value_axis.tick_labels.font.size = Pt(10)


    chart.category_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(8)


    # chart.has_title = True
    # chart.chart_title.text_frame.text = title
    # chart.value_axis.has_major_gridlines = False

    return prs

def plot_heatmap_chart(prs, slide, categories, data_series, title):
    chart_data = XyChartData()
    chart_data.categories = categories
    chart_data = CategoryChartData()

    for series in data_series:
        chart_series = chart_data.add_series(series['name'])
        for value in series['data']:
            chart_series.add_data_point(value)

    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.HEATMAP, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_CHART_TYPE.BOTTOM
    chart.legend.include_in_layout = False

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    return prs